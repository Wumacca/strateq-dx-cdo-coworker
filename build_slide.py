from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colours
STRATEQ_DARK = RGBColor(0x1A, 0x37, 0x5E)   # deep navy
STRATEQ_MID  = RGBColor(0x1F, 0x6F, 0xB2)   # mid blue
STRATEQ_ACC  = RGBColor(0xF0, 0x8C, 0x00)   # amber accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xF4, 0xF6, 0xF9)
MID_GREY     = RGBColor(0xD0, 0xD8, 0xE4)
GREEN        = RGBColor(0x21, 0x96, 0x53)
ORANGE       = RGBColor(0xE6, 0x5C, 0x00)
TEXT_DARK    = RGBColor(0x1A, 0x37, 0x5E)

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width_pt=0.75):
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = Pt(line_width_pt)
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_name="Calibri", font_size=11, bold=False, italic=False,
                 color=TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox, tf

def add_paragraph(tf, text, font_name="Calibri", font_size=10, bold=False, color=TEXT_DARK,
                  align=PP_ALIGN.LEFT, space_before_pt=2, bullet=False, indent_level=0):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before_pt)
    p.level = indent_level
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bullet:
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '▪')
    return p


# ── Slide setup ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

W = prs.slide_width
H = prs.slide_height
M = Inches(0.3)   # margin

# ── Background ───────────────────────────────────────────────────────────────
bg = add_rect(slide, 0, 0, W, H, fill_rgb=LIGHT_GREY)

# ── Header bar ───────────────────────────────────────────────────────────────
HDR_H = Inches(0.68)
add_rect(slide, 0, 0, W, HDR_H, fill_rgb=STRATEQ_DARK)

# Accent line under header
add_rect(slide, 0, HDR_H, W, Pt(3), fill_rgb=STRATEQ_ACC)

add_text_box(slide, "Operations Digital Delivery  –  June 2026",
             M, Inches(0.12), W - 2*M, Inches(0.46),
             font_size=22, bold=True, color=WHITE)

# ── Grid layout constants ─────────────────────────────────────────────────────
Y_TOP    = HDR_H + Inches(0.22)
ROW_H    = Inches(1.05)   # highlight boxes row
BODY_Y   = Y_TOP + ROW_H + Inches(0.15)
BODY_H   = Inches(3.55)
FWD_Y    = BODY_Y + BODY_H + Inches(0.15)
FWD_H    = Inches(1.28)

COL_GAP  = Inches(0.18)
COL1_X   = M
COL1_W   = Inches(2.55)
COL2_X   = COL1_X + COL1_W + COL_GAP
COL2_W   = Inches(4.65)
COL3_X   = COL2_X + COL2_W + COL_GAP
COL3_W   = Inches(4.65)
# (small gap on right is fine)

# ── KPI highlight tiles ───────────────────────────────────────────────────────
tile_data = [
    ("Under Budget", "All initiatives"),
    ("No Additional Cost", "Emergent work absorbed"),
    ("Reg 5 KPIs", "No observations"),
    ("2 Adoption Gates", "COSHH 365 & Omega IM"),
]
n = len(tile_data)
tile_w = (W - 2*M - (n-1)*COL_GAP) / n
for i, (top_line, sub_line) in enumerate(tile_data):
    tx = M + i*(tile_w + COL_GAP)
    add_rect(slide, tx, Y_TOP, tile_w, ROW_H - Inches(0.05),
             fill_rgb=STRATEQ_DARK, line_rgb=None)
    # amber top accent
    add_rect(slide, tx, Y_TOP, tile_w, Pt(4), fill_rgb=STRATEQ_ACC)
    add_text_box(slide, top_line,
                 tx + Inches(0.1), Y_TOP + Inches(0.06),
                 tile_w - Inches(0.2), Inches(0.42),
                 font_size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, sub_line,
                 tx + Inches(0.08), Y_TOP + Inches(0.46),
                 tile_w - Inches(0.16), Inches(0.42),
                 font_size=8.5, color=MID_GREY, align=PP_ALIGN.CENTER)

# ── Section header helper ─────────────────────────────────────────────────────
def section_card(slide, x, y, w, h, title, title_bg=STRATEQ_MID):
    add_rect(slide, x, y, w, h, fill_rgb=WHITE, line_rgb=MID_GREY, line_width_pt=0.5)
    add_rect(slide, x, y, w, Inches(0.32), fill_rgb=title_bg)
    add_text_box(slide, title, x + Inches(0.12), y + Inches(0.04),
                 w - Inches(0.24), Inches(0.26),
                 font_size=10.5, bold=True, color=WHITE)

def body_tf(slide, x, y, w, h, top_pad=Inches(0.36)):
    tb = slide.shapes.add_textbox(x + Inches(0.14), y + top_pad,
                                  w - Inches(0.28), h - top_pad - Inches(0.08))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    return tf

# ── Left card: Digital Initiatives ───────────────────────────────────────────
section_card(slide, COL1_X, BODY_Y, COL2_W, BODY_H, "Digital Initiatives")
tf_di = body_tf(slide, COL1_X, BODY_Y, COL2_W, BODY_H)
tf_di.paragraphs[0].space_before = Pt(0)

items_di = [
    ("KPI Dashboard", "Centralisation & AIRB complete; NEO Next rollout live; automation in progress.", GREEN, False),
    ("Omega 365 — Incident Management", "Live", GREEN, True),
    ("Omega 365 — Action Management", "Live  (Quality rollout delayed)", ORANGE, True),
    ("COSHH 365", "Live; adoption gate approved", GREEN, False),
    ("Salus Suite", "In progress; delayed by Safety Case priorities; training commenced", ORANGE, False),
    ("BP Andrew & NEO Transition Support", "Active", GREEN, False),
    ("Reg 5 Support", "Active", GREEN, False),
    ("Training & Competence", "Vendor selection in progress", ORANGE, False),
    ("CMMS", "Final vendor selection in progress", ORANGE, False),
]

first = True
for label, status, status_col, indent in items_di:
    from pptx.util import Pt as _Pt
    p = tf_di.add_paragraph() if not first else tf_di.paragraphs[0]
    first = False
    p.space_before = _Pt(3 if not indent else 1)
    p.level = 1 if indent else 0
    from pptx.oxml.ns import qn
    from lxml import etree
    pPr = p._p.get_or_add_pPr()
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', '▪' if not indent else '○')

    r1 = p.add_run()
    r1.text = label + ":  "
    r1.font.name = "Calibri"
    r1.font.size = _Pt(9 if not indent else 8.5)
    r1.font.bold = True
    r1.font.color.rgb = TEXT_DARK

    r2 = p.add_run()
    r2.text = status
    r2.font.name = "Calibri"
    r2.font.size = _Pt(9 if not indent else 8.5)
    r2.font.bold = False
    r2.font.color.rgb = status_col

# ── Right card: Chronos Developments ─────────────────────────────────────────
section_card(slide, COL3_X, BODY_Y, COL3_W, BODY_H, "Chronos Developments")
tf_ch = body_tf(slide, COL3_X, BODY_Y, COL3_W, BODY_H)

chron_live = [
    "HR Dashboards",
    "Service Orders / Agreements",
    "New Start Forms",
    "Resource Utilisation",
    "Purchase Requisitions",
]
chron_prog = [
    ("Expenses Enhancements", "Complete; rollout in progress"),
    ("AP Register", "User testing"),
    ("Sales Invoicing", "User testing"),
]
chron_scope = [
    "AP Month-End Automation",
    "Goods Receipt & Manifest",
    "Budget Management (Actual vs Budget)",
]

from pptx.util import Pt as _Pt
from pptx.oxml.ns import qn as _qn
from lxml import etree as _etree

def _bullet_run(tf, text, font_size=9, bold=False, color=TEXT_DARK, first=False, space_before=3):
    p = tf.add_paragraph() if not first else tf.paragraphs[0]
    p.space_before = _Pt(space_before)
    pPr = p._p.get_or_add_pPr()
    bc = _etree.SubElement(pPr, _qn('a:buChar'))
    bc.set('char', '▪')
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = _Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p

def _label_status(tf, label, status, status_col, font_size=9, first=False, space_before=3):
    p = tf.add_paragraph() if not first else tf.paragraphs[0]
    p.space_before = _Pt(space_before)
    pPr = p._p.get_or_add_pPr()
    bc = _etree.SubElement(pPr, _qn('a:buChar'))
    bc.set('char', '▪')
    r1 = p.add_run(); r1.text = label + ":  "
    r1.font.name = "Calibri"; r1.font.size = _Pt(font_size)
    r1.font.bold = True; r1.font.color.rgb = TEXT_DARK
    r2 = p.add_run(); r2.text = status
    r2.font.name = "Calibri"; r2.font.size = _Pt(font_size)
    r2.font.bold = False; r2.font.color.rgb = status_col

# Section sub-header: Live
p0 = tf_ch.paragraphs[0]
p0.space_before = _Pt(0)
r0 = p0.add_run(); r0.text = "Live"
r0.font.name = "Calibri"; r0.font.size = _Pt(9); r0.font.bold = True
r0.font.color.rgb = STRATEQ_MID

for item in chron_live:
    _label_status(tf_ch, item, "Live", GREEN, space_before=2)

# Separator
p_sep = tf_ch.add_paragraph(); p_sep.space_before = _Pt(5)
r_sep = p_sep.add_run(); r_sep.text = "In Progress"
r_sep.font.name = "Calibri"; r_sep.font.size = _Pt(9); r_sep.font.bold = True
r_sep.font.color.rgb = STRATEQ_MID

for label, status in chron_prog:
    _label_status(tf_ch, label, status, ORANGE, space_before=2)

p_sep2 = tf_ch.add_paragraph(); p_sep2.space_before = _Pt(5)
r_sep2 = p_sep2.add_run(); r_sep2.text = "New Development Scoping"
r_sep2.font.name = "Calibri"; r_sep2.font.size = _Pt(9); r_sep2.font.bold = True
r_sep2.font.color.rgb = STRATEQ_MID

for item in chron_scope:
    _bullet_run(tf_ch, item, space_before=2)

# ── Forward View card (bottom, full width) ────────────────────────────────────
section_card(slide, M, FWD_Y, W - 2*M, FWD_H, "Forward View & Pipeline", title_bg=STRATEQ_ACC)

fwd_items = [
    "Complete remaining capitalisation initiatives",
    "Continue BP Andrew transition support",
    "Progress Salus Suite rollout",
    "Finalise CMMS selection",
    "Finalise Training & Competence system selection and roll-out",
    "Prepare next 6-month capitalisation programme",
]

# Two-column layout inside forward view
n2 = len(fwd_items)
half = (n2 + 1) // 2
inner_w = (W - 2*M - Inches(0.28)) / 2
inner_x1 = M + Inches(0.14)
inner_x2 = inner_x1 + inner_w + Inches(0.1)

def fwd_col(slide, x, y, w, h, items):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = _Pt(3)
        pPr = p._p.get_or_add_pPr()
        bc = _etree.SubElement(pPr, _qn('a:buChar'))
        bc.set('char', '▸')
        r = p.add_run(); r.text = item
        r.font.name = "Calibri"; r.font.size = _Pt(9.5)
        r.font.bold = False; r.font.color.rgb = TEXT_DARK

col_y    = FWD_Y + Inches(0.36)
col_h    = FWD_H - Inches(0.44)
fwd_col(slide, inner_x1, col_y, inner_w, col_h, fwd_items[:half])
fwd_col(slide, inner_x2, col_y, inner_w, col_h, fwd_items[half:])

# ── Footer ────────────────────────────────────────────────────────────────────
FOOT_H = Inches(0.28)
add_rect(slide, 0, H - FOOT_H, W, FOOT_H, fill_rgb=STRATEQ_DARK)
add_text_box(slide, "Strateq Digital  |  Confidential  |  June 2026",
             M, H - FOOT_H + Inches(0.04), W/2, Inches(0.22),
             font_size=7.5, color=MID_GREY)
add_text_box(slide, "Operations Digital Delivery",
             W/2, H - FOOT_H + Inches(0.04), W/2 - M, Inches(0.22),
             font_size=7.5, color=MID_GREY, align=PP_ALIGN.RIGHT)

out_path = "/home/user/strateq-dx-cdo-coworker/Operations_Digital_Delivery_June2026.pptx"
prs.save(out_path)
print("Saved:", out_path)
